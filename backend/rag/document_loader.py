import json
import csv
import os
from typing import List, Dict, Any
from pathlib import Path
import pandas as pd
import logging

# File processing imports
import pdfplumber  # pdfplumber for PDF files
from docx import Document  # python-docx for Word documents
from pptx import Presentation  # python-pptx for PowerPoint files

# Setup logging
logger = logging.getLogger(__name__)

class DocumentLoader:
    """Loads and processes documents from the knowledge base"""
    
    def __init__(self, knowledge_base_path: str):
        self.knowledge_base_path = Path(knowledge_base_path)
        self.supported_formats = ['.txt', '.json', '.csv', '.md', '.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']
    
    def load_all_documents(self) -> List[Dict[str, Any]]:
        """Load all documents from the knowledge base"""
        documents = []
        
        for file_path in self.knowledge_base_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_formats:
                try:
                    doc = self.load_document(file_path)
                    if doc:
                        documents.append(doc)
                except Exception as e:
                    logger.error(f"Error loading {file_path}: {e}")
                    print(f"Error loading {file_path}: {e}")
        
        return documents
    
    def load_document(self, file_path: Path) -> Dict[str, Any]:
        """Load a single document based on its file type"""
        suffix = file_path.suffix.lower()
        
        if suffix == '.json':
            return self._load_json(file_path)
        elif suffix == '.csv':
            return self._load_csv(file_path)
        elif suffix == '.txt':
            return self._load_text(file_path)
        elif suffix == '.md':
            return self._load_markdown(file_path)
        elif suffix == '.pdf':
            return self._load_pdf(file_path)
        elif suffix in ['.docx', '.doc']:
            return self._load_word(file_path)
        elif suffix in ['.xlsx', '.xls']:
            return self._load_excel(file_path)
        elif suffix in ['.pptx', '.ppt']:
            return self._load_powerpoint(file_path)
        else:
            raise ValueError(f"Unsupported file format: {suffix}")
    
    def _load_json(self, file_path: Path) -> Dict[str, Any]:
        """Load JSON document"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        # Flatten JSON structure for better processing
        content = self._flatten_json(data)
        
        return {
            'source': str(file_path),
            'title': data.get('title', file_path.stem),
            'category': data.get('category', 'general'),
            'content': content,
            'metadata': {
                'file_type': 'json',
                'file_size': file_path.stat().st_size,
                'last_modified': file_path.stat().st_mtime
            }
        }
    
    def _load_csv(self, file_path: Path) -> Dict[str, Any]:
        """Load CSV document"""
        df = pd.read_csv(file_path)
        
        # Convert DataFrame to structured text
        content_lines = []
        for _, row in df.iterrows():
            line_parts = []
            for col, value in row.items():
                if pd.notna(value):
                    line_parts.append(f"{col}: {value}")
            content_lines.append(" | ".join(line_parts))
        
        content = "\n".join(content_lines)
        
        return {
            'source': str(file_path),
            'title': file_path.stem,
            'category': 'data',
            'content': content,
            'metadata': {
                'file_type': 'csv',
                'columns': list(df.columns),
                'rows': len(df),
                'file_size': file_path.stat().st_size,
                'last_modified': file_path.stat().st_mtime
            }
        }
    
    def _load_text(self, file_path: Path) -> Dict[str, Any]:
        """Load text document"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            'source': str(file_path),
            'title': file_path.stem,
            'category': 'text',
            'content': content,
            'metadata': {
                'file_type': 'txt',
                'file_size': file_path.stat().st_size,
                'last_modified': file_path.stat().st_mtime
            }
        }
    
    def _load_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Load markdown document"""
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        return {
            'source': str(file_path),
            'title': file_path.stem,
            'category': 'markdown',
            'content': content,
            'metadata': {
                'file_type': 'md',
                'file_size': file_path.stat().st_size,
                'last_modified': file_path.stat().st_mtime
            }
        }
    
    def _load_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Load PDF document"""
        try:
            content_lines = []
            page_count = 0
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    text = page.extract_text()
                    if text and text.strip():
                        content_lines.append(f"--- Page {page_num} ---")
                        content_lines.append(text)
                    page_count += 1
            
            content = "\n".join(content_lines)
            
            return {
                'source': str(file_path),
                'title': file_path.stem,
                'category': 'document',
                'content': content,
                'metadata': {
                    'file_type': 'pdf',
                    'pages': page_count,
                    'file_size': file_path.stat().st_size,
                    'last_modified': file_path.stat().st_mtime
                }
            }
        except Exception as e:
            logger.error(f"Error loading PDF {file_path}: {e}")
            raise
    
    def _load_word(self, file_path: Path) -> Dict[str, Any]:
        """Load Word document (.docx)"""
        try:
            doc = Document(file_path)
            content_lines = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    content_lines.append(para.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content_lines.append(" | ".join(row_text))
            
            content = "\n".join(content_lines)
            
            return {
                'source': str(file_path),
                'title': file_path.stem,
                'category': 'document',
                'content': content,
                'metadata': {
                    'file_type': 'docx',
                    'paragraphs': len(doc.paragraphs),
                    'tables': len(doc.tables),
                    'file_size': file_path.stat().st_size,
                    'last_modified': file_path.stat().st_mtime
                }
            }
        except Exception as e:
            logger.error(f"Error loading Word document {file_path}: {e}")
            raise
    
    def _load_excel(self, file_path: Path) -> Dict[str, Any]:
        """Load Excel document"""
        try:
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content_lines = []
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content_lines.append(f"--- Sheet: {sheet_name} ---")
                
                # Add column headers
                if not df.empty:
                    content_lines.append("Columns: " + " | ".join(df.columns.astype(str)))
                    
                    # Add data rows (limit to first 100 rows to avoid huge content)
                    for idx, row in df.head(100).iterrows():
                        row_data = []
                        for col, value in row.items():
                            if pd.notna(value):
                                row_data.append(f"{col}: {value}")
                        if row_data:
                            content_lines.append(" | ".join(row_data))
                    
                    if len(df) > 100:
                        content_lines.append(f"... and {len(df) - 100} more rows")
            
            content = "\n".join(content_lines)
            
            return {
                'source': str(file_path),
                'title': file_path.stem,
                'category': 'spreadsheet',
                'content': content,
                'metadata': {
                    'file_type': 'xlsx',
                    'sheets': excel_file.sheet_names,
                    'total_sheets': len(excel_file.sheet_names),
                    'file_size': file_path.stat().st_size,
                    'last_modified': file_path.stat().st_mtime
                }
            }
        except Exception as e:
            logger.error(f"Error loading Excel file {file_path}: {e}")
            raise
    
    def _load_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Load PowerPoint document"""
        try:
            prs = Presentation(file_path)
            content_lines = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content_lines.append(f"--- Slide {slide_num} ---")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_lines.append(shape.text)
                
                # Extract text from tables
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                content_lines.append(" | ".join(row_text))
            
            content = "\n".join(content_lines)
            
            return {
                'source': str(file_path),
                'title': file_path.stem,
                'category': 'presentation',
                'content': content,
                'metadata': {
                    'file_type': 'pptx',
                    'slides': len(prs.slides),
                    'file_size': file_path.stat().st_size,
                    'last_modified': file_path.stat().st_mtime
                }
            }
        except Exception as e:
            logger.error(f"Error loading PowerPoint file {file_path}: {e}")
            raise

    def _flatten_json(self, obj: Any, prefix: str = '') -> str:
        """Flatten JSON object into readable text optimized for LLM understanding"""
        if isinstance(obj, dict):
            lines = []
            for key, value in obj.items():
                if isinstance(value, list):
                    # Handle lists specially for better readability
                    lines.append(f"{key}:")
                    for i, item in enumerate(value):
                        if isinstance(item, dict):
                            lines.append(f"  {i+1}. {self._format_dict_as_text(item)}")
                        else:
                            lines.append(f"  {i+1}. {item}")
                elif isinstance(value, dict):
                    lines.append(f"{key}:")
                    lines.append(self._format_dict_as_text(value, indent="  "))
                else:
                    lines.append(f"{key}: {value}")
            return "\n".join(lines)
        elif isinstance(obj, list):
            lines = []
            for i, item in enumerate(obj):
                if isinstance(item, dict):
                    lines.append(f"Item {i+1}: {self._format_dict_as_text(item)}")
                else:
                    lines.append(f"Item {i+1}: {item}")
            return "\n".join(lines)
        else:
            return str(obj)
    
    def _format_dict_as_text(self, obj: dict, indent: str = "") -> str:
        """Format a dictionary as readable text"""
        lines = []
        for key, value in obj.items():
            if isinstance(value, dict):
                lines.append(f"{indent}{key}:")
                lines.append(self._format_dict_as_text(value, indent + "  "))
            elif isinstance(value, list):
                lines.append(f"{indent}{key}:")
                for i, item in enumerate(value):
                    if isinstance(item, dict):
                        lines.append(f"{indent}  {i+1}. {self._format_dict_as_text(item, indent + '    ')}")
                    else:
                        lines.append(f"{indent}  {i+1}. {item}")
            else:
                lines.append(f"{indent}{key}: {value}")
        return "\n".join(lines)
